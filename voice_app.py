import os
import json
import base64
import io
import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from upstash_redis import Redis

# Librerie per la lettura dei vari tipi di file
from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation

# 1. Configurazione della pagina
st.set_page_config(
    page_title="Fixi",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# 2. Caricamento variabili d'ambiente
load_dotenv()

api_key = st.secrets.get("API_KEY") or os.getenv("API_KEY")
redis_url = st.secrets.get("UPSTASH_REDIS_REST_URL") or os.getenv("UPSTASH_REDIS_REST_URL")
redis_token = st.secrets.get("UPSTASH_REDIS_REST_TOKEN") or os.getenv("UPSTASH_REDIS_REST_TOKEN")

client = OpenAI(
    api_key=api_key,
    base_url="https://api.groq.com/openai/v1"
)

redis = Redis(url=redis_url, token=redis_token) if (redis_url and redis_token) else None

# 3. Funzione Universale per Processare QUALSIASI File
def process_uploaded_file(uploaded_file, openai_client):
    file_name = uploaded_file.name
    file_bytes = uploaded_file.getvalue()
    file_ext = os.path.splitext(file_name)[1].lower()

    # --- IMMAGINI ---
    if file_ext in [".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"]:
        base64_img = base64.b64encode(file_bytes).decode('utf-8')
        mime = uploaded_file.type or f"image/{file_ext.replace('.', '')}"
        return {
            "is_image": True,
            "base64": base64_img,
            "mime": mime,
            "name": file_name,
            "info": f"📸 Immagine: **{file_name}**"
        }

    # --- AUDIO (Whisper) ---
    elif file_ext in [".mp3", ".wav", ".m4a", ".ogg", ".webm", ".flac"]:
        try:
            audio_file = (file_name, file_bytes, uploaded_file.type or "audio/mpeg")
            transcription = openai_client.audio.transcriptions.create(
                model="whisper-large-v3-turbo",
                file=audio_file
            )
            return {"is_image": False, "text": f"\n\n[TRASCRIZIONE AUDIO '{file_name}']:\n{transcription.text}", "info": f"🎙️ Audio **{file_name}** trascritto!"}
        except Exception as e:
            return {"is_image": False, "text": f"\n\n[AUDIO '{file_name}' - Errore: {e}]", "info": f"⚠️ Errore audio: {file_name}"}

    # --- PDF ---
    elif file_ext == ".pdf":
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text = "\n".join([page.extract_text() or "" for page in reader.pages]).strip()
            return {"is_image": False, "text": f"\n\n[CONTENUTO PDF '{file_name}']:\n{text or '[PDF privo di testo]'}", "info": f"📄 PDF **{file_name}**"}
        except Exception as e:
            return {"is_image": False, "text": f"\n\n[PDF '{file_name}']: Errore ({e})", "info": f"⚠️ Errore PDF: {file_name}"}

    # --- WORD (.docx) ---
    elif file_ext in [".docx", ".doc"]:
        try:
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            return {"is_image": False, "text": f"\n\n[CONTENUTO WORD '{file_name}']:\n{text}", "info": f"📝 Word **{file_name}**"}
        except Exception as e:
            return {"is_image": False, "text": f"\n\n[WORD '{file_name}']: Errore ({e})", "info": f"⚠️ Errore Word: {file_name}"}

    # --- EXCEL / CSV ---
    elif file_ext in [".xlsx", ".xls", ".csv"]:
        try:
            df = pd.read_csv(io.BytesIO(file_bytes)) if file_ext == ".csv" else pd.read_excel(io.BytesIO(file_bytes))
            return {"is_image": False, "text": f"\n\n[TABELLA '{file_name}']:\n{df.head(100).to_markdown(index=False)}", "info": f"📊 Tabella **{file_name}**"}
        except Exception as e:
            return {"is_image": False, "text": f"\n\n[TABELLA '{file_name}']: Errore ({e})", "info": f"⚠️ Errore tabella: {file_name}"}

    # --- POWERPOINT (.pptx) ---
    elif file_ext == ".pptx":
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if slide_text:
                    text_runs.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
            return {"is_image": False, "text": f"\n\n[POWERPOINT '{file_name}']:\n" + "\n\n".join(text_runs), "info": f"📊 PowerPoint **{file_name}**"}
        except Exception as e:
            return {"is_image": False, "text": f"\n\n[POWERPOINT '{file_name}']: Errore ({e})", "info": f"⚠️ Errore PPT: {file_name}"}

    # --- CODICE E TESTO (LUA, HTML, CSS, JS, PY, JSON, XML, TXT, ecc.) ---
    try:
        text = file_bytes.decode("utf-8")
        return {"is_image": False, "text": f"\n\n[FILE CODICE/TESTO '{file_name}']:\n{text}", "info": f"💻 File **{file_name}**"}
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
            return {"is_image": False, "text": f"\n\n[FILE CODICE/TESTO '{file_name}']:\n{text}", "info": f"💻 File **{file_name}**"}
        except Exception:
            pass

    # --- BINARI GENERICI ---
    size_kb = round(len(file_bytes) / 1024, 2)
    return {
        "is_image": False,
        "text": f"\n\n[ALLEGATO BINARIO - Nome: '{file_name}', Dimensione: {size_kb} KB]",
        "info": f"📦 File **{file_name}** ({size_kb} KB)"
    }

# 4. Conoscenza da GitHub
def load_github_knowledge():
    knowledge = ""
    allowed_extensions = ('.txt', '.md', '.csv', '.json')
    ignored_files = {'requirements.txt'}

    try:
        for file in os.listdir("."):
            if file.endswith(allowed_extensions) and file not in ignored_files:
                try:
                    with open(file, "r", encoding="utf-8") as f:
                        knowledge += f"\n--- CONTENUTO FILE '{file}' ---\n{f.read()}\n"
                except Exception:
                    pass
    except Exception as e:
        print(f"Errore lettura file: {e}")
        
    return knowledge

# 5. Pulizia dello Storico (rimuove il base64 pesante per le richieste successive)
def sanitize_messages_for_api(messages):
    clean_messages = []
    for msg in messages:
        role = msg["role"]
        content = msg["content"]
        
        if isinstance(content, list):
            # Se è una lista con immagini allegate
            text_parts = []
            has_image = False
            for part in content:
                if isinstance(part, dict):
                    if part.get("type") == "text":
                        text_parts.append(part.get("text", ""))
                    elif part.get("type") == "image_url":
                        has_image = True
            
            combined = " ".join(text_parts)
            if has_image and "[Immagine allegata]" not in combined:
                combined += " [Immagine allegata precedente]"
            
            clean_messages.append({"role": role, "content": combined.strip() or "[Allegato]"})
        else:
            clean_messages.append({"role": role, "content": str(content)})
            
    return clean_messages

# 6. Gestione Memoria Cloud
def load_chat_history():
    if redis:
        try:
            data = redis.get("fixi_chat_history")
            if data:
                return json.loads(data)
        except Exception:
            pass
    return []

def save_chat_history(messages):
    if redis:
        try:
            clean_msgs = sanitize_messages_for_api(messages)
            redis.set("fixi_chat_history", json.dumps(clean_msgs, ensure_ascii=False))
        except Exception as e:
            st.error(f"Errore salvataggio memoria: {e}")

if "messages" not in st.session_state:
    st.session_state.messages = load_chat_history()

# 7. Stili CSS
st.markdown("""
<style>
    * { box-sizing: border-box; }
    #MainMenu, header, footer, .stDeployButton { display: none !important; }
    
    div[data-testid="stAppViewContainer"], div[data-testid="stMain"] {
        background-color: #0d0f14 !important;
    }

    body { 
        margin: 0; padding: 0; 
        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; 
        background-color: #0d0f14; color: #e0e0e0; 
    }

    .block-container {
        padding: 20px 20px 180px 20px !important; 
        max-width: 900px !important;
        margin: 0 auto !important;
    }

    .glow-purple {
        position: fixed; bottom: -10%; left: -5%; width: 500px; height: 500px;
        background: radial-gradient(circle, rgba(70, 45, 140, 0.2) 0%, rgba(13, 15, 20, 0) 70%);
        pointer-events: none; z-index: 1;
    }
    .glow-cyan {
        position: fixed; top: -10%; right: -5%; width: 550px; height: 550px;
        background: radial-gradient(circle, rgba(30, 90, 140, 0.15) 0%, rgba(13, 15, 20, 0) 70%);
        pointer-events: none; z-index: 1;
    }

    .header-section {
        display: flex; flex-direction: column; align-items: center; margin-bottom: 25px; z-index: 5;
    }
    .waveform-container {
        position: relative; width: 75px; height: 75px; display: flex; align-items: center; justify-content: center; margin-bottom: 4px;
    }
    .wave-ring-1 {
        position: absolute; width: 100%; height: 100%;
        border-radius: 42% 58% 60% 40% / 45% 55% 45% 55%;
        border: 2px solid #38bdf8;
        box-shadow: 0 0 10px rgba(56, 189, 248, 0.6), inset 0 0 6px rgba(56, 189, 248, 0.2);
        animation: waveMorph1 6s ease-in-out infinite alternate;
    }
    .wave-ring-2 {
        position: absolute; width: 95%; height: 95%;
        border-radius: 55% 45% 38% 62% / 50% 42% 58% 50%;
        border: 2px solid #c084fc;
        box-shadow: 0 0 10px rgba(192, 132, 252, 0.6), inset 0 0 6px rgba(192, 132, 252, 0.2);
        animation: waveMorph2 7s ease-in-out infinite alternate;
    }
    .wave-ring-3 {
        position: absolute; width: 90%; height: 90%;
        border-radius: 48% 52% 55% 45% / 60% 40% 60% 40%;
        border: 2px solid #818cf8;
        box-shadow: 0 0 8px rgba(129, 140, 248, 0.4);
        animation: waveMorph3 5s ease-in-out infinite alternate;
    }

    @keyframes waveMorph1 {
        0% { transform: rotate(0deg) scale(0.97); }
        100% { transform: rotate(360deg) scale(0.97); }
    }
    @keyframes waveMorph2 {
        0% { transform: rotate(0deg) scale(1.01); }
        100% { transform: rotate(-360deg) scale(1.01); }
    }
    @keyframes waveMorph3 {
        0% { transform: rotate(0deg) scale(0.98); }
        100% { transform: rotate(180deg) scale(1.03); }
    }
    .status-text { font-size: 11px; color: rgba(255, 255, 255, 0.6); }

    div[data-testid="stChatMessage"] {
        background: rgba(255, 255, 255, 0.015);
        border: 1px solid rgba(255, 255, 255, 0.05);
        border-radius: 12px;
        padding: 10px 15px;
        margin-bottom: 15px;
        z-index: 5;
    }

    div[data-testid="stChatInput"] {
        position: fixed !important;
        bottom: 25px !important;
        left: 50% !important;
        transform: translateX(-50%) !important;
        width: calc(100% - 40px) !important;
        max-width: 900px !important;
        z-index: 100 !important;
        background: transparent !important;
    }
    div[data-testid="stChatInput"] textarea {
        background-color: rgba(13, 15, 20, 0.95) !important;
        color: white !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        border-radius: 12px !important;
        box-shadow: 0 -10px 25px rgba(13, 15, 20, 0.9) !important;
    }
</style>

<div class="glow-purple"></div>
<div class="glow-cyan"></div>
""", unsafe_allow_html=True)

# Header
st.markdown("""
<div class="header-section">
    <div class="waveform-container">
        <div class="wave-ring-1"></div>
        <div class="wave-ring-2"></div>
        <div class="wave-ring-3"></div>
    </div>
    <div class="status-text">Fixi ti sta ascoltando...</div>
</div>
""", unsafe_allow_html=True)

# Barra con pulsante di pulizia chat
col_clear, _ = st.columns([1, 4])
with col_clear:
    if st.button("🗑️ Pulisci conversazione", use_container_width=True):
        st.session_state.messages = []
        if redis:
            try:
                redis.delete("fixi_chat_history")
            except Exception:
                pass
        st.rerun()

github_docs = load_github_knowledge()

system_prompt = {
    "role": "system",
    "content": f"Sei Fixi, un'intelligenza artificiale esperta di programmazione (in particolare FiveM, Lua, JS, HTML/CSS). Rispondi in modo chiaro, preciso e strutturato in italiano.\n"
               f"File di contesto del proprietario:\n{github_docs}"
}

# Mostra i messaggi precedenti
for message in st.session_state.messages:
    if message["role"] != "system":
        with st.chat_message(message["role"]):
            if isinstance(message["content"], str):
                st.write(message["content"])
            elif isinstance(message["content"], list):
                for part in message["content"]:
                    if part.get("type") == "text":
                        st.write(part["text"])
                    elif part.get("type") == "image_url":
                        st.image(part["image_url"]["url"], width=280)

# Pulsante Tasto "+" per allegare file
with st.popover("➕ Allega file", use_container_width=True):
    uploaded_files = st.file_uploader(
        "Scegli file dal dispositivo", 
        type=None, 
        accept_multiple_files=True
    )

combined_text = ""
first_image = None

if uploaded_files:
    for u_file in uploaded_files:
        res = process_uploaded_file(u_file, client)
        st.write(res["info"])
        if res.get("is_image") and not first_image:
            first_image = res
        elif not res.get("is_image"):
            combined_text += res.get("text", "")

# Input della Chat
prompt = st.chat_input("Scrivi un messaggio a Fixi...")

if prompt:
    user_text = str(prompt)
    
    # Prepara il messaggio utente corrente
    if first_image:
        model_to_use = "llama-3.2-11b-vision-preview"
        full_text = (user_text + combined_text) if (user_text or combined_text) else "Analizza questo file/immagine."
        current_msg = {
            "role": "user",
            "content": [
                {"type": "text", "text": full_text},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{first_image['mime']};base64,{first_image['base64']}"}
                }
            ]
        }
    else:
        model_to_use = "llama-3.3-70b-versatile"
        current_msg = {"role": "user", "content": user_text + combined_text}

    # Prepara gli storici passati puliti
    past_clean_messages = sanitize_messages_for_api(st.session_state.messages)
    
    # Costruisci la lista dei messaggi da inviare all'API
    api_payload = [system_prompt] + past_clean_messages + [current_msg]

    # Mostra subito a schermo
    st.session_state.messages.append(current_msg)
    with st.chat_message("user"):
        st.write(user_text if user_text else "Inviato allegato")

    with st.spinner("Fixi sta elaborando..."):
        try:
            risposta = client.chat.completions.create(
                model=model_to_use,
                messages=api_payload
            )
            testo_risposta = risposta.choices[0].message.content
            
            st.session_state.messages.append({"role": "assistant", "content": testo_risposta})
            save_chat_history(st.session_state.messages)
            st.rerun()
            
        except Exception as e:
            st.error(f"⚠️ **Errore API:** {e}\n\n*Consiglio: Prova a cliccare su '🗑️ Pulisci conversazione' in alto se il problema persiste.*")
